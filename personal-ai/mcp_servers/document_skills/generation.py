"""Document Skills MCP 使用的受限格式生成器。"""

from __future__ import annotations

import re
from html import escape
from pathlib import Path
from typing import Any, Callable

from core.files.artifacts import (
    ArtifactError,
    allocate_artifact,
    artifact_tool_result,
    discard_artifact,
    finalize_artifact,
    resolve_artifact,
)


MAX_TEXT_CHARS = 100_000
MAX_SLIDES = 20
MAX_SHEETS = 10
MAX_ROWS = 2_000
MAX_COLUMNS = 50
ACCENT = "2563EB"
DARK = "172033"


def _require_text(value: str, label: str, maximum: int = MAX_TEXT_CHARS) -> str:
    text = str(value or "").strip()
    if not text:
        raise ValueError(f"{label}不能为空")
    if len(text) > maximum:
        raise ValueError(f"{label}超过长度限制")
    return text


def _blocks(content: str) -> list[tuple[str, str, int]]:
    content = _require_text(content, "内容")
    result: list[tuple[str, str, int]] = []
    paragraph: list[str] = []

    def flush() -> None:
        if paragraph:
            result.append(("paragraph", " ".join(paragraph).strip(), 0))
            paragraph.clear()

    for raw in content.splitlines():
        line = raw.strip()
        if not line:
            flush()
            continue
        heading = re.match(r"^(#{1,3})\s+(.+)$", line)
        bullet = re.match(r"^[-*]\s+(.+)$", line)
        numbered = re.match(r"^\d+[.)]\s+(.+)$", line)
        if heading:
            flush()
            result.append(("heading", heading.group(2).strip(), len(heading.group(1))))
        elif bullet:
            flush()
            result.append(("bullet", bullet.group(1).strip(), 0))
        elif numbered:
            flush()
            result.append(("number", numbered.group(1).strip(), 0))
        else:
            paragraph.append(line)
    flush()
    return result


def _create(
    filename: str,
    suffix: str,
    builder: Callable[[Path], None],
) -> str:
    artifact_id, path, display_name = allocate_artifact(filename, suffix)
    try:
        builder(path)
        record = finalize_artifact(artifact_id, path, display_name)
        return artifact_tool_result(record)
    except Exception:
        discard_artifact(artifact_id)
        raise


def create_docx(title: str, content: str, filename: str = "document.docx") -> str:
    title = _require_text(title, "标题", 200)
    blocks = _blocks(content)

    def build(path: Path) -> None:
        try:
            from docx import Document
            from docx.enum.section import WD_SECTION
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.oxml import OxmlElement
            from docx.oxml.ns import qn
            from docx.shared import Cm, Pt, RGBColor
        except ImportError as exc:
            raise ValueError("缺少 python-docx 依赖") from exc

        document = Document()
        section = document.sections[0]
        section.page_width, section.page_height = Cm(21), Cm(29.7)
        section.top_margin = section.bottom_margin = Cm(2.4)
        section.left_margin = section.right_margin = Cm(2.5)

        normal = document.styles["Normal"]
        normal.font.name = "Aptos"
        normal.font.size = Pt(11)
        normal._element.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")
        normal.paragraph_format.space_after = Pt(7)
        normal.paragraph_format.line_spacing = 1.25
        for level, size in ((1, 18), (2, 15), (3, 13)):
            style = document.styles[f"Heading {level}"]
            style.font.name = "Aptos Display"
            style.font.size = Pt(size)
            style.font.bold = True
            style.font.color.rgb = RGBColor.from_string(DARK if level > 1 else ACCENT)
            style._element.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")
            style.paragraph_format.space_before = Pt(14)
            style.paragraph_format.space_after = Pt(6)

        title_paragraph = document.add_paragraph()
        title_paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
        title_paragraph.paragraph_format.space_after = Pt(18)
        run = title_paragraph.add_run(title)
        run.bold = True
        run.font.name = "Aptos Display"
        run.font.size = Pt(26)
        run.font.color.rgb = RGBColor.from_string(DARK)
        run._element.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")

        for kind, text, level in blocks:
            if kind == "heading":
                document.add_paragraph(text, style=f"Heading {level}")
            elif kind == "bullet":
                document.add_paragraph(text, style="List Bullet")
            elif kind == "number":
                document.add_paragraph(text, style="List Number")
            else:
                document.add_paragraph(text)

        footer = section.footer.paragraphs[0]
        footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
        footer_run = footer.add_run("Personal AI · Document Skills")
        footer_run.font.size = Pt(8)
        footer_run.font.color.rgb = RGBColor(120, 126, 138)
        document.core_properties.title = title
        document.core_properties.author = "Personal AI"
        document.save(path)

    return _create(filename, ".docx", build)


def append_docx(artifact_id: str, content: str) -> str:
    record = resolve_artifact(artifact_id)
    if record.path.suffix.lower() != ".docx":
        raise ArtifactError("只能编辑由 Document Skills 生成的 DOCX")
    blocks = _blocks(content)
    try:
        from docx import Document
    except ImportError as exc:
        raise ValueError("缺少 python-docx 依赖") from exc
    document = Document(record.path)
    for kind, text, level in blocks:
        if kind == "heading":
            document.add_paragraph(text, style=f"Heading {level}")
        elif kind == "bullet":
            document.add_paragraph(text, style="List Bullet")
        elif kind == "number":
            document.add_paragraph(text, style="List Number")
        else:
            document.add_paragraph(text)
    document.save(record.path)
    updated = finalize_artifact(record.id, record.path, record.filename)
    return artifact_tool_result(updated)


def create_pdf(title: str, content: str, filename: str = "document.pdf") -> str:
    title = _require_text(title, "标题", 200)
    blocks = _blocks(content)

    def build(path: Path) -> None:
        try:
            from reportlab.lib import colors
            from reportlab.lib.enums import TA_LEFT
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
            from reportlab.lib.units import mm
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.ttfonts import TTFont
            from reportlab.platypus import (
                ListFlowable, ListItem, Paragraph, SimpleDocTemplate, Spacer,
            )
        except ImportError as exc:
            raise ValueError("缺少 reportlab 依赖") from exc

        font_path = Path("C:/Windows/Fonts/simhei.ttf")
        font_name = "ArtifactCJK"
        if font_path.is_file() and font_name not in pdfmetrics.getRegisteredFontNames():
            pdfmetrics.registerFont(TTFont(font_name, str(font_path)))
        if font_name not in pdfmetrics.getRegisteredFontNames():
            font_name = "Helvetica"

        styles = getSampleStyleSheet()
        body = ParagraphStyle(
            "ArtifactBody", parent=styles["BodyText"], fontName=font_name,
            fontSize=10.5, leading=17, textColor=colors.HexColor(f"#{DARK}"),
            spaceAfter=7, wordWrap="CJK",
        )
        title_style = ParagraphStyle(
            "ArtifactTitle", parent=body, fontSize=24, leading=31,
            textColor=colors.HexColor(f"#{DARK}"), spaceAfter=16,
        )
        heading_styles = {
            level: ParagraphStyle(
                f"ArtifactH{level}", parent=body, fontSize=size, leading=size + 7,
                textColor=colors.HexColor(f"#{ACCENT if level == 1 else DARK}"),
                spaceBefore=10, spaceAfter=5,
            )
            for level, size in ((1, 17), (2, 14), (3, 12))
        }
        story: list[Any] = [Paragraph(escape(title), title_style)]
        pending_list: list[Any] = []
        pending_kind = ""

        def flush_list() -> None:
            nonlocal pending_list, pending_kind
            if pending_list:
                story.append(ListFlowable(
                    pending_list,
                    bulletType="1" if pending_kind == "number" else "bullet",
                    leftIndent=18,
                    bulletFontName=font_name,
                ))
                story.append(Spacer(1, 4))
                pending_list = []
                pending_kind = ""

        for kind, text, level in blocks:
            if kind in {"bullet", "number"}:
                if pending_kind and pending_kind != kind:
                    flush_list()
                pending_kind = kind
                pending_list.append(ListItem(Paragraph(escape(text), body)))
                continue
            flush_list()
            story.append(Paragraph(escape(text), heading_styles[level] if kind == "heading" else body))
        flush_list()

        def page_footer(canvas, document) -> None:
            canvas.saveState()
            canvas.setFont(font_name, 8)
            canvas.setFillColor(colors.HexColor("#7A8190"))
            canvas.drawCentredString(A4[0] / 2, 12 * mm, f"{document.page}")
            canvas.restoreState()

        document = SimpleDocTemplate(
            str(path), pagesize=A4, rightMargin=22 * mm, leftMargin=22 * mm,
            topMargin=22 * mm, bottomMargin=22 * mm,
            title=title, author="Personal AI",
        )
        document.build(story, onFirstPage=page_footer, onLaterPages=page_footer)

    return _create(filename, ".pdf", build)


def create_pptx(title: str, slides: list[dict], filename: str = "presentation.pptx") -> str:
    title = _require_text(title, "标题", 200)
    if not isinstance(slides, list) or not 1 <= len(slides) <= MAX_SLIDES:
        raise ValueError(f"幻灯片数量必须在 1 到 {MAX_SLIDES} 之间")
    normalized: list[tuple[str, list[str]]] = []
    for index, slide in enumerate(slides, 1):
        if not isinstance(slide, dict):
            raise ValueError(f"第 {index} 页必须是对象")
        slide_title = _require_text(slide.get("title", ""), f"第 {index} 页标题", 160)
        bullets = slide.get("bullets", [])
        if not isinstance(bullets, list) or len(bullets) > 10:
            raise ValueError("每页 bullets 必须是不超过 10 项的数组")
        normalized.append((slide_title, [_require_text(item, "要点", 600) for item in bullets]))

    def build(path: Path) -> None:
        try:
            from pptx import Presentation
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            from pptx.util import Inches, Pt
        except ImportError as exc:
            raise ValueError("缺少 python-pptx 依赖") from exc

        deck = Presentation()
        deck.slide_width, deck.slide_height = Inches(13.333), Inches(7.5)
        cover = deck.slides.add_slide(deck.slide_layouts[6])
        background = cover.background.fill
        background.solid(); background.fore_color.rgb = RGBColor.from_string(DARK)
        box = cover.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(11.5), Inches(2.2))
        paragraph = box.text_frame.paragraphs[0]
        paragraph.text = title
        paragraph.font.name = "Microsoft YaHei"
        paragraph.font.size = Pt(40)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.LEFT
        subtitle = cover.shapes.add_textbox(Inches(0.95), Inches(4.65), Inches(7), Inches(0.5))
        sub_p = subtitle.text_frame.paragraphs[0]
        sub_p.text = "Personal AI · Document Skills"
        sub_p.font.name = "Microsoft YaHei"; sub_p.font.size = Pt(16)
        sub_p.font.color.rgb = RGBColor(170, 190, 225)

        for number, (slide_title, bullets) in enumerate(normalized, 1):
            slide = deck.slides.add_slide(deck.slide_layouts[6])
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(248, 250, 252)
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(11.7), Inches(0.8))
            title_p = title_box.text_frame.paragraphs[0]
            title_p.text = slide_title
            title_p.font.name = "Microsoft YaHei"; title_p.font.size = Pt(30)
            title_p.font.bold = True; title_p.font.color.rgb = RGBColor.from_string(DARK)
            accent = slide.shapes.add_shape(1, Inches(0.8), Inches(1.5), Inches(0.9), Inches(0.08))
            accent.fill.solid(); accent.fill.fore_color.rgb = RGBColor.from_string(ACCENT)
            accent.line.fill.background()
            body_box = slide.shapes.add_textbox(Inches(1), Inches(1.9), Inches(11.2), Inches(4.7))
            frame = body_box.text_frame; frame.word_wrap = True; frame.clear()
            for item_index, item in enumerate(bullets or ["待补充"]):
                item_p = frame.paragraphs[0] if item_index == 0 else frame.add_paragraph()
                item_p.text = item; item_p.level = 0
                item_p.font.name = "Microsoft YaHei"; item_p.font.size = Pt(20)
                item_p.font.color.rgb = RGBColor.from_string(DARK)
                item_p.space_after = Pt(14)
                item_p.text = f"•  {item}"
            page_box = slide.shapes.add_textbox(Inches(12.1), Inches(6.9), Inches(0.5), Inches(0.3))
            page_p = page_box.text_frame.paragraphs[0]
            page_p.text = str(number); page_p.font.size = Pt(10)
            page_p.font.color.rgb = RGBColor(120, 126, 138)
        deck.core_properties.title = title
        deck.core_properties.author = "Personal AI"
        deck.save(path)

    return _create(filename, ".pptx", build)


def create_xlsx(title: str, sheets: list[dict], filename: str = "workbook.xlsx") -> str:
    title = _require_text(title, "标题", 200)
    if not isinstance(sheets, list) or not 1 <= len(sheets) <= MAX_SHEETS:
        raise ValueError(f"工作表数量必须在 1 到 {MAX_SHEETS} 之间")

    def build(path: Path) -> None:
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Alignment, Font, PatternFill
            from openpyxl.utils import get_column_letter
        except ImportError as exc:
            raise ValueError("缺少 openpyxl 依赖") from exc

        workbook = Workbook()
        workbook.remove(workbook.active)
        used_names: set[str] = set()
        for index, spec in enumerate(sheets, 1):
            if not isinstance(spec, dict):
                raise ValueError(f"第 {index} 个工作表必须是对象")
            name = _require_text(spec.get("name", f"Sheet{index}"), "工作表名称", 31)
            name = re.sub(r"[\\/*?:\[\]]", "_", name)
            if name in used_names:
                raise ValueError(f"工作表名称重复：{name}")
            used_names.add(name)
            rows = spec.get("rows", [])
            if not isinstance(rows, list) or len(rows) > MAX_ROWS:
                raise ValueError(f"工作表 {name} 行数超过限制")
            sheet = workbook.create_sheet(name)
            for row_index, row in enumerate(rows, 1):
                if not isinstance(row, list) or len(row) > MAX_COLUMNS:
                    raise ValueError(f"工作表 {name} 第 {row_index} 行格式错误")
                for col_index, value in enumerate(row, 1):
                    if value is not None and not isinstance(value, (str, int, float, bool)):
                        raise ValueError("单元格只能是文本、数字、布尔值或空值")
                    sheet.cell(row_index, col_index, value)
            sheet.freeze_panes = "A2" if rows else None
            sheet.auto_filter.ref = sheet.dimensions if rows else None
            if rows:
                for cell in sheet[1]:
                    cell.font = Font(name="Microsoft YaHei", bold=True, color="FFFFFF")
                    cell.fill = PatternFill("solid", fgColor=ACCENT)
                    cell.alignment = Alignment(vertical="center", wrap_text=True)
                sheet.row_dimensions[1].height = 25
            for row in sheet.iter_rows(min_row=2):
                for cell in row:
                    cell.font = Font(name="Microsoft YaHei", size=10)
                    cell.alignment = Alignment(vertical="top", wrap_text=True)
            for col_index in range(1, min(MAX_COLUMNS, sheet.max_column) + 1):
                values = [str(sheet.cell(row, col_index).value or "") for row in range(1, min(sheet.max_row, 200) + 1)]
                width = min(40, max(10, max((len(value) for value in values), default=8) + 2))
                sheet.column_dimensions[get_column_letter(col_index)].width = width
        workbook.properties.title = title
        workbook.properties.creator = "Personal AI"
        workbook.save(path)

    return _create(filename, ".xlsx", build)
